from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = ROOT / "export"
OUTPUT = ROOT / "getplaced-hackathon-deck.pptx"


def main() -> None:
    slide_images = sorted(EXPORT_DIR.glob("slide-*.png"))
    if len(slide_images) != 12:
        raise RuntimeError(f"Expected 12 slide images, found {len(slide_images)}")

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    blank_layout = deck.slide_layouts[6]

    deck.core_properties.title = "getPlaced Hackathon Presentation"
    deck.core_properties.subject = "Product showcase"
    deck.core_properties.author = "TheMentalists"

    for slide_image in slide_images:
        slide = deck.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(slide_image),
            0,
            0,
            width=deck.slide_width,
            height=deck.slide_height,
        )

    deck.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(deck.slides)} slides")


if __name__ == "__main__":
    main()
